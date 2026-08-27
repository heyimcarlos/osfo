import argparse
import base64
import io
import json
import math
import os
import re
import subprocess
import tempfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


MAX_SLIDES = 20
MAX_VISUAL_EDGE = 2048
TITLE_FONT = "Liberation Sans"
BODY_FONT = "Liberation Sans"


def read_source(path: Path) -> dict[str, object]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError("artifact source must be an object")
    return value


def decode_visuals(value: object) -> dict[str, bytes]:
    if not isinstance(value, list):
        raise ValueError("supporting visuals must be a list")
    decoded: dict[str, bytes] = {}
    total = 0
    for item in value:
        if not isinstance(item, dict):
            raise ValueError("supporting visual must be an object")
        content_id = item.get("contentId")
        encoded = item.get("base64")
        if not isinstance(content_id, str) or not isinstance(encoded, str):
            raise ValueError("supporting visual identity and body are required")
        body = base64.b64decode(encoded, validate=True)
        total += len(body)
        if total > 25_000_000:
            raise ValueError("supporting visuals exceed the immutable input limit")
        with Image.open(io.BytesIO(body)) as image:
            image.verify()
        decoded[content_id] = body
    return decoded


def font(size: int) -> ImageFont.FreeTypeFont:
    candidates = (
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    )
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    raise ValueError("required artifact font is missing")


def render_diagram(source: dict[str, object], output: Path) -> dict[str, object]:
    width = source.get("width")
    height = source.get("height")
    nodes = source.get("nodes")
    edges = source.get("edges")
    direction = source.get("direction")
    if (
        not isinstance(width, int)
        or not isinstance(height, int)
        or not 1 <= width <= MAX_VISUAL_EDGE
        or not 1 <= height <= MAX_VISUAL_EDGE
        or not isinstance(nodes, list)
        or not 1 <= len(nodes) <= 20
        or not isinstance(edges, list)
        or len(edges) > 40
        or direction not in ("leftToRight", "topToBottom")
    ):
        raise ValueError("diagram exceeds its bounded semantic source")
    image = Image.new("RGBA", (width, height), "white")
    draw = ImageDraw.Draw(image)
    title_font = font(max(18, min(36, width // 24)))
    label_font = font(max(14, min(26, width // 34)))
    title = source.get("title")
    if not isinstance(title, str) or not title:
        raise ValueError("diagram title is required")
    draw.text((width // 2, 24), title, font=title_font, fill="#172554", anchor="ma")
    columns = len(nodes) if direction == "leftToRight" else min(3, len(nodes))
    rows = math.ceil(len(nodes) / columns)
    margin_x = max(28, width // 20)
    top = max(90, height // 6)
    available_height = height - top - 30
    cell_width = (width - 2 * margin_x) / columns
    cell_height = available_height / rows
    box_width = min(cell_width * 0.72, 300)
    box_height = min(cell_height * 0.55, 120)
    positions: dict[str, tuple[float, float, float, float]] = {}
    for index, node in enumerate(nodes):
        if not isinstance(node, dict) or not isinstance(node.get("id"), str):
            raise ValueError("diagram node is invalid")
        column = index if direction == "leftToRight" else index % columns
        row = 0 if direction == "leftToRight" else index // columns
        center_x = margin_x + cell_width * (column + 0.5)
        center_y = top + cell_height * (row + 0.5)
        positions[node["id"]] = (
            center_x - box_width / 2,
            center_y - box_height / 2,
            center_x + box_width / 2,
            center_y + box_height / 2,
        )
    for edge in edges:
        if not isinstance(edge, dict):
            raise ValueError("diagram edge is invalid")
        start = positions.get(edge.get("from"))
        finish = positions.get(edge.get("to"))
        if start is None or finish is None:
            raise ValueError("diagram edge references an unknown node")
        start_point = ((start[0] + start[2]) / 2, (start[1] + start[3]) / 2)
        finish_point = ((finish[0] + finish[2]) / 2, (finish[1] + finish[3]) / 2)
        draw.line((start_point, finish_point), fill="#64748b", width=4)
        edge_label = edge.get("label")
        if not isinstance(edge_label, str):
            raise ValueError("diagram edge label is invalid")
        if edge_label:
            label_center = (
                (start_point[0] + finish_point[0]) / 2,
                (start_point[1] + finish_point[1]) / 2,
            )
            bounds = draw.textbbox(label_center, edge_label, font=label_font, anchor="mm")
            draw.rounded_rectangle(
                (bounds[0] - 6, bounds[1] - 3, bounds[2] + 6, bounds[3] + 3),
                radius=4,
                fill="white",
            )
            draw.text(label_center, edge_label, font=label_font, fill="#334155", anchor="mm")
    for node in nodes:
        box = positions[node["id"]]
        draw.rounded_rectangle(box, radius=18, fill="#dbeafe", outline="#1d4ed8", width=3)
        label = node.get("label")
        if not isinstance(label, str) or not label:
            raise ValueError("diagram node label is required")
        measured = draw.textbbox((0, 0), label, font=label_font)
        if measured[2] - measured[0] > box_width - 24:
            raise ValueError("diagram label does not fit its node")
        draw.text(
            ((box[0] + box[2]) / 2, (box[1] + box[3]) / 2),
            label,
            font=label_font,
            fill="#172554",
            anchor="mm",
        )
    image.save(output, format="PNG", optimize=True)
    return {"height": height, "kind": "visual", "width": width}


def normalize_image(source: dict[str, object], provider_body: str, output: Path) -> dict[str, object]:
    width = source.get("width")
    height = source.get("height")
    if not isinstance(width, int) or not isinstance(height, int):
        raise ValueError("image dimensions are required")
    body = base64.b64decode(provider_body, validate=True)
    with Image.open(io.BytesIO(body)) as image:
        converted = image.convert("RGBA")
        if converted.size != (width, height):
            converted = converted.resize((width, height), Image.Resampling.LANCZOS)
        converted.save(output, format="PNG", optimize=True)
    return {"height": height, "kind": "visual", "width": width}


def clear_slides(presentation: Presentation) -> None:
    slide_ids = presentation.slides._sldIdLst  # type: ignore[attr-defined]
    for slide_id in list(slide_ids):
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        slide_ids.remove(slide_id)


def add_notes(slide, notes: str, sources: list[str]) -> None:
    lines = [notes] if notes else []
    if sources:
        lines.extend(["[Sources]", *sources])
    if not lines:
        return
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "\n".join(lines)


def inspect_rendered_presentation(output: Path, expected_slides: int) -> list[str]:
    """Render the actual PPTX, then inspect every full-size rasterized slide."""
    issues: list[str] = []
    with tempfile.TemporaryDirectory(prefix="osfo-pptx-inspection-") as directory:
        root = Path(directory)
        environment = {**os.environ, "HOME": directory, "SAL_USE_VCLPLUGIN": "svp"}
        converted = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                directory,
                str(output),
            ],
            check=False,
            capture_output=True,
            env=environment,
            text=True,
            timeout=30,
        )
        pdf = root / f"{output.stem}.pdf"
        if converted.returncode != 0 or not pdf.is_file():
            raise ValueError(
                "generated PPTX could not be rendered: "
                f"stdout={converted.stdout!r} stderr={converted.stderr!r}"
            )
        rasterized = subprocess.run(
            ["pdftoppm", "-png", "-r", "120", str(pdf), str(root / "slide")],
            check=False,
            capture_output=True,
            text=True,
            timeout=30,
        )
        if rasterized.returncode != 0:
            raise ValueError("rendered PPTX pages could not be rasterized")
        pages = sorted(root.glob("slide-*.png"))
        if len(pages) != expected_slides:
            raise ValueError("actual rendered slide count does not match the presentation")
        for index, page in enumerate(pages, start=1):
            with Image.open(page) as image:
                image.load()
                rendered = image.convert("RGB")
                if rendered.width < 1200 or rendered.height < 675:
                    issues.append(f"slide {index}: rendered below full-size inspection resolution")
                extrema = rendered.convert("L").getextrema()
                if extrema is None or extrema[1] - extrema[0] < 24:
                    issues.append(f"slide {index}: unreadable or blank rendered output")
    return issues


def inspect_text_fit(item: dict[str, object], has_visual: bool) -> list[str]:
    title = item.get("title")
    body = item.get("body")
    if not isinstance(title, str) or not isinstance(body, list):
        raise ValueError("presentation text is invalid")
    title_font = font(60)
    body_font = font(37)
    issues: list[str] = []
    if title_font.getlength(title) > 1416:
        issues.append("title wraps in its fixed text box")
    body_width = 852 if has_visual else 1416
    if any(not isinstance(line, str) or body_font.getlength(line) > body_width for line in body):
        issues.append("body wraps in its fixed text box")
    if len(body) * 60 > 630:
        issues.append("body clips in its fixed text box")
    return issues


def render_presentation(
    source: dict[str, object],
    visuals: dict[str, bytes],
    output: Path,
    source_presentation: Path | None,
) -> dict[str, object]:
    slides = source.get("slides")
    if not isinstance(slides, list) or not 1 <= len(slides) <= MAX_SLIDES:
        raise ValueError("presentation must contain between 1 and 20 slides")
    presentation = Presentation(str(source_presentation)) if source_presentation else Presentation()
    if source_presentation:
        clear_slides(presentation)
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    layout = presentation.slide_layouts[6]
    issues: list[str] = []
    for index, item in enumerate(slides):
        if not isinstance(item, dict):
            raise ValueError("presentation slide is invalid")
        title = item.get("title")
        body = item.get("body")
        if not isinstance(title, str) or not title or not isinstance(body, list):
            raise ValueError("presentation slide title and body are required")
        if "\n" in title or "\r" in title:
            raise ValueError("presentation titles must be one line")
        slide = presentation.slides.add_slide(layout)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 250, 252)
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = False
        title_run = title_frame.paragraphs[0].add_run()
        title_run.text = title
        title_run.font.name = TITLE_FONT
        title_run.font.size = Pt(36)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(23, 37, 84)
        image_id = item.get("imageContentId")
        diagram_id = item.get("diagramContentId")
        visual_id = image_id if isinstance(image_id, str) else diagram_id
        body_width = 7.1 if isinstance(visual_id, str) else 11.8
        body_box = slide.shapes.add_textbox(
            Inches(0.85), Inches(1.55), Inches(body_width), Inches(5.25)
        )
        body_frame = body_box.text_frame
        body_frame.clear()
        body_frame.word_wrap = True
        for body_index, line in enumerate(body):
            if not isinstance(line, str):
                raise ValueError("presentation body line is invalid")
            if "\n" in line or "\r" in line:
                raise ValueError("presentation body items must be one line")
            paragraph = body_frame.paragraphs[0] if body_index == 0 else body_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.name = BODY_FONT
            paragraph.font.size = Pt(22)
            paragraph.font.color.rgb = RGBColor(30, 41, 59)
            paragraph.space_after = Pt(14)
        if isinstance(visual_id, str):
            visual = visuals.get(visual_id)
            if visual is None:
                raise ValueError("presentation references an unavailable visual")
            with Image.open(io.BytesIO(visual)) as image:
                image.verify()
            slide.shapes.add_picture(io.BytesIO(visual), Inches(8.35), Inches(1.65), width=Inches(4.2))
        else:
            visual = None
        sources = item.get("sourceNotes")
        notes = item.get("speakerNotes")
        if not isinstance(sources, list) or any(not isinstance(value, str) for value in sources):
            raise ValueError("presentation sources are invalid")
        add_notes(slide, notes if isinstance(notes, str) else "", sources)
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                issues.append(f"slide {index + 1}: shape outside canvas")
            if shape.left + shape.width > presentation.slide_width:
                issues.append(f"slide {index + 1}: horizontal overflow")
            if shape.top + shape.height > presentation.slide_height:
                issues.append(f"slide {index + 1}: vertical overflow")
        if len(title) > 70:
            issues.append(f"slide {index + 1}: title may wrap")
        if len(body) > 10 or any(len(str(line)) > 140 for line in body):
            issues.append(f"slide {index + 1}: body may overflow")
        issues.extend(
            f"slide {index + 1}: {issue}"
            for issue in inspect_text_fit(item, isinstance(visual_id, str))
        )
    presentation.save(output)
    issues.extend(inspect_rendered_presentation(output, len(slides)))
    rendered_count = len(slides)
    return {"issues": issues, "kind": "presentation", "renderedSlideCount": rendered_count}


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--kind", choices=("presentation", "image", "diagram"), required=True)
    parser.add_argument("--input", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--source-presentation", type=Path)
    arguments = parser.parse_args()
    source_envelope = read_source(arguments.input)
    source = source_envelope.get("source")
    if not isinstance(source, dict):
        raise ValueError("artifact source is missing")
    if arguments.kind == "presentation":
        inspection = render_presentation(
            source,
            decode_visuals(source_envelope.get("supportingVisuals")),
            arguments.output,
            arguments.source_presentation,
        )
    elif arguments.kind == "image":
        provider_image = source_envelope.get("providerImageBase64")
        if not isinstance(provider_image, str):
            raise ValueError("verified provider image is missing")
        inspection = normalize_image(source, provider_image, arguments.output)
    else:
        inspection = render_diagram(source, arguments.output)
    print(json.dumps(inspection, separators=(",", ":")))


if __name__ == "__main__":
    main()
